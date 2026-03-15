"""Nexus knowledge service — file parsing, chunking, AI summarization."""

import json
import logging
import os
from pathlib import Path

from database.connection import get_connection, row_to_dict, rows_to_dicts
from services.ai_provider import generate_ai_response, check_ai_available

logger = logging.getLogger(__name__)

UPLOADS_DIR = Path(__file__).resolve().parent.parent.parent / "backend" / "uploads"

# ---------------------------------------------------------------------------
# Parse queue management
# ---------------------------------------------------------------------------


def enqueue_parse(file_id: int) -> dict:
    """Insert a file into the parse queue."""
    with get_connection() as conn:
        with conn.cursor() as cur:
            cur.execute(
                """INSERT INTO knowledge_parse_queue (file_id, status)
                   VALUES (%s, 'pending')
                   RETURNING *""",
                (file_id,),
            )
            return row_to_dict(cur)


def get_next_parse_task() -> dict | None:
    """Get next pending task with row-level lock (skip locked)."""
    with get_connection() as conn:
        with conn.cursor() as cur:
            cur.execute(
                """SELECT * FROM knowledge_parse_queue
                   WHERE status = 'pending'
                   ORDER BY created_at ASC
                   LIMIT 1
                   FOR UPDATE SKIP LOCKED"""
            )
            task = row_to_dict(cur)
            if task:
                cur.execute(
                    """UPDATE knowledge_parse_queue
                       SET status = 'processing'
                       WHERE id = %s""",
                    (task["id"],),
                )
            return task


def update_parse_status(
    queue_id: int, status: str, error_message: str | None = None
) -> None:
    """Update queue task status."""
    with get_connection() as conn:
        with conn.cursor() as cur:
            cur.execute(
                """UPDATE knowledge_parse_queue
                   SET status = %s, error_message = %s, processed_at = NOW()
                   WHERE id = %s""",
                (status, error_message, queue_id),
            )


# ---------------------------------------------------------------------------
# Knowledge CRUD
# ---------------------------------------------------------------------------


def create_knowledge_chunk(
    file_id: int,
    client_id: int | None,
    chunk_index: int,
    content: str,
    summary: str | None = None,
    tags: list[str] | None = None,
    metadata: dict | None = None,
) -> dict:
    with get_connection() as conn:
        with conn.cursor() as cur:
            cur.execute(
                """INSERT INTO nx_knowledge
                   (file_id, client_id, chunk_index, content, summary, tags, metadata)
                   VALUES (%s, %s, %s, %s, %s, %s, %s)
                   RETURNING *""",
                (
                    file_id,
                    client_id,
                    chunk_index,
                    content,
                    summary,
                    tags or [],
                    json.dumps(metadata) if metadata else None,
                ),
            )
            return row_to_dict(cur)


def get_knowledge_by_file(file_id: int) -> list[dict]:
    """Get all knowledge chunks for a file, ordered by chunk_index."""
    with get_connection() as conn:
        with conn.cursor() as cur:
            cur.execute(
                """SELECT * FROM nx_knowledge
                   WHERE file_id = %s
                   ORDER BY chunk_index""",
                (file_id,),
            )
            return rows_to_dicts(cur)


def get_knowledge_by_client(client_id: int) -> list[dict]:
    """Get all knowledge chunks across files for a client."""
    with get_connection() as conn:
        with conn.cursor() as cur:
            cur.execute(
                """SELECT k.*, f.file_name
                   FROM nx_knowledge k
                   JOIN nx_file f ON k.file_id = f.id
                   WHERE k.client_id = %s
                   ORDER BY k.file_id, k.chunk_index""",
                (client_id,),
            )
            return rows_to_dicts(cur)


def delete_knowledge_by_file(file_id: int) -> int:
    """Delete all knowledge chunks for a file (before re-parse)."""
    with get_connection() as conn:
        with conn.cursor() as cur:
            cur.execute(
                "DELETE FROM nx_knowledge WHERE file_id = %s",
                (file_id,),
            )
            return cur.rowcount


def list_all_knowledge(limit: int = 100) -> list[dict]:
    """Get all knowledge chunks across all files."""
    with get_connection() as conn:
        with conn.cursor() as cur:
            cur.execute(
                """SELECT k.*, f.file_name
                   FROM nx_knowledge k
                   JOIN nx_file f ON k.file_id = f.id
                   ORDER BY k.file_id, k.chunk_index
                   LIMIT %s""",
                (limit,),
            )
            return rows_to_dicts(cur)


def search_knowledge(
    query: str, client_id: int | None = None
) -> list[dict]:
    """Search knowledge chunks by ILIKE on content + summary."""
    like_pattern = f"%{query}%"
    with get_connection() as conn:
        with conn.cursor() as cur:
            if client_id:
                cur.execute(
                    """SELECT k.*, f.file_name
                       FROM nx_knowledge k
                       JOIN nx_file f ON k.file_id = f.id
                       WHERE k.client_id = %s
                         AND (k.content ILIKE %s OR k.summary ILIKE %s)
                       ORDER BY k.file_id, k.chunk_index
                       LIMIT 50""",
                    (client_id, like_pattern, like_pattern),
                )
            else:
                cur.execute(
                    """SELECT k.*, f.file_name
                       FROM nx_knowledge k
                       JOIN nx_file f ON k.file_id = f.id
                       WHERE k.content ILIKE %s OR k.summary ILIKE %s
                       ORDER BY k.file_id, k.chunk_index
                       LIMIT 50""",
                    (like_pattern, like_pattern),
                )
            return rows_to_dicts(cur)


def update_file_parse_status(file_id: int, status: str) -> None:
    """Update nx_file.parse_status."""
    with get_connection() as conn:
        with conn.cursor() as cur:
            cur.execute(
                "UPDATE nx_file SET parse_status = %s WHERE id = %s",
                (status, file_id),
            )


# ---------------------------------------------------------------------------
# Text extraction
# ---------------------------------------------------------------------------

PARSEABLE_EXTENSIONS = {".pdf", ".pptx", ".docx"}


def _get_file_extension(file_name: str) -> str:
    return os.path.splitext(file_name)[1].lower()


def extract_text_from_file(
    file_path: str, file_name: str
) -> list[dict]:
    """Extract text chunks from a file.

    Returns list of {content: str, metadata: dict}.
    """
    ext = _get_file_extension(file_name)
    full_path = UPLOADS_DIR / os.path.basename(file_path)

    if not full_path.exists():
        raise FileNotFoundError(f"File not found: {full_path}")

    if ext == ".pdf":
        return _extract_pdf(full_path)
    if ext == ".pptx":
        return _extract_pptx(full_path)
    if ext == ".docx":
        return _extract_docx(full_path)

    raise ValueError(f"Unsupported file type: {ext}")


def _extract_pdf(path: Path) -> list[dict]:
    """Extract text from PDF — one chunk per page, split long pages."""
    import pdfplumber

    chunks = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            text = text.strip()
            if not text:
                continue
            if len(text) > 1000:
                # Split on double newline
                parts = text.split("\n\n")
                buffer = ""
                for part in parts:
                    if buffer and len(buffer) + len(part) > 1000:
                        chunks.append({
                            "content": buffer.strip(),
                            "metadata": {"page_number": i + 1},
                        })
                        buffer = part
                    else:
                        buffer = f"{buffer}\n\n{part}" if buffer else part
                if buffer.strip():
                    chunks.append({
                        "content": buffer.strip(),
                        "metadata": {"page_number": i + 1},
                    })
            else:
                chunks.append({
                    "content": text,
                    "metadata": {"page_number": i + 1},
                })
    return chunks


def _extract_pptx(path: Path) -> list[dict]:
    """Extract text from PPTX — one chunk per slide."""
    from pptx import Presentation

    prs = Presentation(str(path))
    chunks = []
    for i, slide in enumerate(prs.slides):
        texts = []
        slide_title = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    texts.append(text)
            if shape.has_text_frame and shape.shape_id == slide.shapes.title_shape_id if hasattr(slide.shapes, 'title_shape_id') else False:
                slide_title = shape.text_frame.text.strip()
        # Try to get title from slide.shapes.title
        if not slide_title and slide.shapes.title:
            slide_title = slide.shapes.title.text.strip() if slide.shapes.title.has_text_frame else None
        content = "\n".join(texts)
        if not content.strip():
            continue
        chunks.append({
            "content": content,
            "metadata": {
                "slide_number": i + 1,
                "slide_title": slide_title,
            },
        })
    return chunks


def _extract_docx(path: Path) -> list[dict]:
    """Extract text from DOCX — accumulate paragraphs into 500-1000 char chunks."""
    from docx import Document

    doc = Document(str(path))
    chunks = []
    buffer = ""
    current_section = None

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        is_heading = para.style and para.style.name.startswith("Heading")

        if is_heading:
            # Flush buffer before new section
            if buffer.strip():
                chunks.append({
                    "content": buffer.strip(),
                    "metadata": {"section_title": current_section},
                })
                buffer = ""
            current_section = text
            buffer = text
        elif len(buffer) + len(text) > 1000:
            # Flush when buffer exceeds threshold
            if buffer.strip():
                chunks.append({
                    "content": buffer.strip(),
                    "metadata": {"section_title": current_section},
                })
            buffer = text
        else:
            buffer = f"{buffer}\n{text}" if buffer else text

    # Flush remaining
    if buffer.strip():
        chunks.append({
            "content": buffer.strip(),
            "metadata": {"section_title": current_section},
        })

    return chunks


# ---------------------------------------------------------------------------
# AI summarization
# ---------------------------------------------------------------------------

_SUMMARIZE_PROMPT = """You are a knowledge extraction assistant for a B2B sales CRM system.
Given a text chunk from a business document, provide:
1. A concise summary (1-3 sentences, in the same language as the content)
2. 3-8 relevant tags for categorization (in the same language as the content)

Respond in JSON format only:
{"summary": "...", "tags": ["tag1", "tag2", ...]}"""


def summarize_chunk(content: str) -> dict:
    """Use AI to generate summary and tags for a chunk.

    Returns {"summary": str | None, "tags": list[str]}.
    Falls back gracefully on failure.
    """
    available, _ = check_ai_available()
    if not available:
        return {"summary": None, "tags": []}

    try:
        # Truncate very long content
        truncated = content[:3000] if len(content) > 3000 else content
        response = generate_ai_response(_SUMMARIZE_PROMPT, truncated)

        # Try to parse JSON from response
        # Strip markdown code fences if present
        cleaned = response.strip()
        if cleaned.startswith("```"):
            lines = cleaned.split("\n")
            cleaned = "\n".join(lines[1:])
            if cleaned.endswith("```"):
                cleaned = cleaned[:-3]
            cleaned = cleaned.strip()

        parsed = json.loads(cleaned)
        return {
            "summary": parsed.get("summary"),
            "tags": parsed.get("tags", []),
        }
    except Exception as e:
        logger.warning("AI summarize failed: %s", e)
        return {"summary": None, "tags": []}
